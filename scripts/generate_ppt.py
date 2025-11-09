"""
Generate a 10-slide PowerPoint for Chessy using images in slides_images/ and slide text.

Usage:
    python scripts\generate_ppt.py

If `python-pptx` is not installed, the script will inform you how to install it.
"""
import os
import sys

OUT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', 'slides_images'))
OUT_PPTX = os.path.join(OUT_DIR, 'Chessy_presentation.pptx')

slides_images = {
    'top_moves_hist': os.path.join(OUT_DIR, 'top_moves_hist.png'),
    'piece_heatmap': os.path.join(OUT_DIR, 'piece_heatmap.png'),
    'feature_importance': os.path.join(OUT_DIR, 'feature_importance.png'),
    'move_history_example': os.path.join(OUT_DIR, 'move_history_example.png'),
    'flow_diagram': os.path.join(OUT_DIR, 'flow_diagram.png'),
    'architecture_diagram': os.path.join(OUT_DIR, 'architecture_diagram.png'),
}

slide_texts = {
    1: {
        'title': 'Chessy — RandomForest Chess Engine with GUI',
        'bullets': ['Lightweight chess engine + Tkinter GUI', 'Team: [Your Name] (Maker) | Guide: [Advisor Name]', r'Demo + Source: c:\\Users\\91738\\OneDrive\\Desktop\\coding\\Chessy']
    },
    2: {
        'title': 'Problem Statement & Objective',
        'bullets': [
            'Problem: Casual players lack an interpretable, fast local engine',
            'Objective: RF-based move prediction + minimax with custom eval',
            'Constraint: Runs locally, minimal deps, responsive GUI'
        ]
    },
    3: {
        'title': 'Dataset Description',
        'bullets': [
            'Source: PGN historical games (data/sample.pgn)',
            'Granularity: positions -> moves (UCI strings)',
            'Features: board encoding + engineered features',
            'Target: next move index (data/vocab.json)'
        ]
    },
    4: {
        'title': 'Data Cleaning & EDA',
        'bullets': [
            'Parse PGN -> positions; build vocab (vocab.py)',
            'Encode board; compute features (train.extract_features)',
            'Saved dataset as NPZ (preprocess.py)'
        ]
    },
    5: {'title': 'Model Selection & Justification', 'bullets': ['RandomForest (scikit-learn): interpretable, fast', 'Hybrid: RF priors + minimax search for legality & lookahead']},
    6: {'title': 'Hyperparameter Tuning & Evaluation', 'bullets': ['Tuning: manual/grid-like (max_depth, min_samples_leaf, n_estimators)', 'Example: max_depth=12, min_samples_leaf=10, n_estimators=30', 'Report: test accuracy / qualitative gameplay']},
    7: {'title': 'Deployment & Architecture', 'bullets': ['Local desktop app: Python + Tkinter GUI', 'Core files: vocab.py, preprocess.py, train.py, play.py, chess_gui.py']},
    8: {'title': 'GUI Demo / Screenshot', 'bullets': ['Drag & drop moves; text-entry; SAN move history', 'Bot thinking indicator and status bar']},
    9: {'title': 'Results & Performance Metrics', 'bullets': ['Training time, test accuracy (~X%)', 'Average bot move time by phase (e.g., depth=3 -> ~1-3s)']},
    10: {'title': 'Conclusion & Future Scope', 'bullets': ['Achievements: local engine + GUI', 'Future: PGN export, time controls, stronger eval, headless tests']}
}

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
except Exception:
    print('\nThe package `python-pptx` is required to generate the PPTX.')
    print('Install with: pip install python-pptx')
    sys.exit(1)

prs = Presentation()
# set slide size if desired (default is fine)

# helper to add title slide
layout_title = prs.slide_layouts[0]
slide = prs.slides.add_slide(layout_title)
title = slide.shapes.title
subtitle = slide.placeholders[1]
slide1 = slide_texts[1]
title.text = slide1['title']
subtitle.text = ' | '.join(slide1['bullets'][0:2])
# Add footer line as small textbox
left = Inches(0.2)
width = Inches(9)
top = Inches(6.8)
height = Inches(0.3)
fol = slide.shapes.add_textbox(left, top, width, height)
f = fol.text_frame
p = f.paragraphs[0]
p.text = slide1['bullets'][2]
p.font.size = Pt(10)

# For slides 2..10: title + bullets + optional image
for i in range(2, 11):
    layout = prs.slide_layouts[1]  # title and content
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide_texts[i]['title']
    # content placeholder
    body = s.shapes.placeholders[1].text_frame
    body.clear()
    for b in slide_texts[i]['bullets']:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)

    # attach images for certain slides
    # slide 3: small table -> attach top_moves_hist if present
    if i == 3 and os.path.exists(slides_images['top_moves_hist']):
        img_path = slides_images['top_moves_hist']
        left = Inches(5.2)
        top = Inches(1.5)
        height = Inches(3.5)
        s.shapes.add_picture(img_path, left, top, height=height)
    if i == 4 and os.path.exists(slides_images['piece_heatmap']):
        img_path = slides_images['piece_heatmap']
        left = Inches(5.2)
        top = Inches(1.5)
        height = Inches(3.5)
        s.shapes.add_picture(img_path, left, top, height=height)
    if i == 5 and os.path.exists(slides_images['feature_importance']):
        img_path = slides_images['feature_importance']
        left = Inches(5.2)
        top = Inches(1.2)
        height = Inches(3.5)
        s.shapes.add_picture(img_path, left, top, height=height)
    if i == 8 and os.path.exists(slides_images['move_history_example']):
        img_path = slides_images['move_history_example']
        left = Inches(1)
        top = Inches(2.8)
        height = Inches(3)
        s.shapes.add_picture(img_path, left, top, height=height)
    if i == 6 and os.path.exists(slides_images['flow_diagram']):
        img_path = slides_images['flow_diagram']
        left = Inches(4.8)
        top = Inches(1.2)
        height = Inches(3)
        s.shapes.add_picture(img_path, left, top, height=height)
    if i == 7 and os.path.exists(slides_images['architecture_diagram']):
        img_path = slides_images['architecture_diagram']
        left = Inches(4.8)
        top = Inches(1.2)
        height = Inches(3)
        s.shapes.add_picture(img_path, left, top, height=height)

# Save
os.makedirs(OUT_DIR, exist_ok=True)
prs.save(OUT_PPTX)
print(f'PPTX generated: {OUT_PPTX}')
